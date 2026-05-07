"""
PPT数字人视频生成服务
功能：上传PPT -> 数字人实时讲解 -> 屏幕录制 -> 保存视频
"""

import os
import json
import uuid
import asyncio
import subprocess
import base64
import shutil
from datetime import datetime
from pathlib import Path
from fastapi import FastAPI, UploadFile, File, HTTPException, Request
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from io import BytesIO

# 阿里云通义千问配置
DASHSCOPE_API_KEY = "sk-a546f070b73d43459f3bca60e7ecf140"

app = FastAPI(title="PPT数字人视频制作")

# 配置路径
BASE_DIR = Path(__file__).parent
UPLOAD_DIR = BASE_DIR / "uploads"
PPT_IMAGES_DIR = BASE_DIR / "ppt_images"
TEMP_DIR = BASE_DIR / "temp"

for d in [UPLOAD_DIR, PPT_IMAGES_DIR, TEMP_DIR]:
    d.mkdir(exist_ok=True)

# 先挂载API路由（通过FastAPI的路由系统自动处理）


# 然后挂载静态文件服务，确保API路由优先
@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    return JSONResponse(
        status_code=exc.status_code,
        content={"code": exc.status_code, "message": exc.detail},
    )


@app.exception_handler(Exception)
async def unhandled_exception_handler(request: Request, exc: Exception):
    import traceback

    traceback.print_exc()
    return JSONResponse(
        status_code=500,
        content={"code": 500, "message": str(exc) or "Internal Server Error"},
    )


app.mount("/ppt_images", StaticFiles(directory=str(PPT_IMAGES_DIR)), name="ppt_images")
app.mount("/temp", StaticFiles(directory=str(TEMP_DIR)), name="temp")

# 根路径的静态文件服务，处理HTML文件
from fastapi.responses import FileResponse

@app.get("/")
async def root():
    return FileResponse(str(BASE_DIR / "index.html"), media_type="text/html")

@app.get("/index.html")
async def index():
    return FileResponse(str(BASE_DIR / "index.html"), media_type="text/html")

@app.get("/video-preview.html")
async def video_preview():
    return FileResponse(str(BASE_DIR / "video-preview.html"), media_type="text/html")

# PPT解析和图片转换工具
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from PIL import Image, ImageDraw, ImageFont
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("警告: python-pptx 或 pillow 未安装，PPT解析功能受限")

import requests

WINDOWS_TTS_VOICE_HINTS = {
    "longxia": ["Microsoft Xiaoxiao Desktop", "Microsoft Huihui Desktop"],
    "longwan": ["Microsoft Xiaoyi Desktop", "Microsoft Xiaoxiao Desktop"],
    "longmale": ["Microsoft Yunxi Desktop", "Microsoft Kangkang Desktop"],
    "longnv": ["Microsoft Xiaoxiao Desktop", "Microsoft Huihui Desktop"],
    "longlao": ["Microsoft Huihui Desktop", "Microsoft Xiaoyi Desktop"],
    "longyu": ["Microsoft Xiaoxiao Desktop", "Microsoft Xiaoyi Desktop"],
}

def call_llm_generate_script(slide_content, slide_num):
    """调用LLM为PPT内容生成讲解脚本"""
    try:
        url = "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {DASHSCOPE_API_KEY}",
            "Content-Type": "application/json"
        }

        prompt = f"""你是一个老师，请为以下PPT内容生成一段讲解文字（约60-100字），要自然流畅，像在上课讲解：

PPT第{slide_num}页内容：
{slide_content}

要求：
1. 讲解要通俗易懂
2. 不要太长，60-100字即可
3. 语气要像老师上课
4. 直接返回讲解文字，不要加引号或其他标记"""

        payload = {
            "model": "qwen-plus",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.7
        }

        response = requests.post(url, headers=headers, json=payload, timeout=30)
        if response.status_code == 200:
            result = response.json()
            if "choices" in result and len(result["choices"]) > 0:
                script = result["choices"][0]["message"]["content"]
                return script.strip()
        return None
    except Exception as e:
        print(f"LLM生成脚本错误: {e}")
        return None

def call_tts(text, output_path, voice_id="longxia"):
    """调用阿里云DashScope TTS"""
    try:
        # 使用正确的 TTS API 端点
        tts_url = "https://dashscope.aliyuncs.com/api/v1/services/audio/tts/generation"
        headers = {
            "Authorization": f"Bearer {DASHSCOPE_API_KEY}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": "cosyvoice-tts-1.0",
            "input": {"text": text},
            "parameters": {
                "voice": voice_id,
                "format": "mp3",
                "sample_rate": 16000,
                "speed": 1.0,
                "volume": 1.0,
                "pitch": 1.0
            }
        }

        response = requests.post(tts_url, headers=headers, json=payload, timeout=120)
        print(f"TTS请求: {tts_url}")
        print(f"TTS参数: {payload}")
        print(f"TTS响应状态: {response.status_code}")
        print(f"TTS响应内容: {response.text[:500]}")
        
        if response.status_code == 200:
            result = response.json()
            if "output" in result and "audio_url" in result["output"]:
                audio_url = result["output"]["audio_url"]
                print(f"获取音频URL: {audio_url}")
                audio_resp = requests.get(audio_url, timeout=60)
                if audio_resp.status_code == 200:
                    with open(output_path, 'wb') as f:
                        f.write(audio_resp.content)
                    file_size = os.path.getsize(output_path) if os.path.exists(output_path) else 0
                    if file_size < 1024:
                        print(f"Audio file too small, treat as invalid: {output_path}, size={file_size}")
                        return False
                    print(f"音频保存成功: {output_path}")
                    return True
        print(f"TTS响应: {response.status_code} - {response.text[:500]}")
        return False
    except Exception as e:
        print(f"TTS错误: {e}")
        import traceback
        traceback.print_exc()
        return False


def generate_dummy_wav(output_path, duration_seconds=3, sample_rate=16000):
    """Generate a playable silence WAV for fallback when TTS fails."""
    try:
        import wave
        import struct

        duration_seconds = max(1, int(duration_seconds))
        num_samples = duration_seconds * sample_rate

        with wave.open(output_path, "wb") as wf:
            wf.setnchannels(1)
            wf.setsampwidth(2)  # 16-bit PCM
            wf.setframerate(sample_rate)
            silence_frame = struct.pack("<h", 0)
            wf.writeframes(silence_frame * num_samples)

        file_size = os.path.getsize(output_path) if os.path.exists(output_path) else 0
        print(f"Generated fallback WAV: {output_path}, size={file_size}")
        return True
    except Exception as e:
        print(f"Failed to generate fallback WAV: {e}")
        return False

def generate_dummy_audio(output_path):
    """生成一个简单的音频文件作为备用方案"""
    try:
        # 创建一个简单的静音MP3文件
        import os
        # 写入一个最小的MP3文件头
        with open(output_path, 'wb') as f:
            # MP3文件头
            f.write(b'ID3\x03\x00\x00\x00\x00\x00\x00\x00\x00')
            # 简单的音频数据
            f.write(b'\xFF\xFB\x18\x00\x00\x00\x00')
        print(f"生成备用音频文件: {output_path}")
        return True
    except Exception as e:
        print(f"生成备用音频失败: {e}")
        return False

def generate_dummy_audio(output_path):
    """Generate a playable fallback audio file.

    The video preview uses `<audio>` in browser; writing a "fake mp3 header"
    is not reliably playable. We generate a short silence WAV instead.
    """
    try:
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        return generate_dummy_wav(output_path)
    except Exception as e:
        print(f"Failed to generate fallback audio: {e}")
        return False

def _audio_file_is_valid(file_path, min_size=1024):
    try:
        return os.path.exists(file_path) and os.path.getsize(file_path) >= min_size
    except OSError:
        return False

def _save_audio_bytes(output_path, content):
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    with open(output_path, "wb") as f:
        f.write(content)
    return _audio_file_is_valid(output_path)

def _extract_audio_url(result):
    if not isinstance(result, dict):
        return None

    output = result.get("output")
    if isinstance(output, dict):
        if output.get("audio_url"):
            return output["audio_url"]
        audio = output.get("audio")
        if isinstance(audio, dict) and audio.get("url"):
            return audio["url"]

    data = result.get("data")
    if isinstance(data, dict):
        if data.get("audio_url"):
            return data["audio_url"]
        audio = data.get("audio")
        if isinstance(audio, dict) and audio.get("url"):
            return audio["url"]

    return None

def _powershell_quote(value):
    return "'" + str(value).replace("'", "''") + "'"

def generate_windows_tts_audio(text, output_path, voice_id="longxia"):
    """Use Windows built-in TTS as a spoken fallback when online TTS fails."""
    text = " ".join(str(text or "").split())
    if not text:
        return False

    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    voice_candidates = WINDOWS_TTS_VOICE_HINTS.get(voice_id, [])
    quoted_candidates = ", ".join(_powershell_quote(name) for name in voice_candidates)
    ps_voice_array = f"@({quoted_candidates})" if quoted_candidates else "@()"

    ps_script = f"""
$ErrorActionPreference = 'Stop'
Add-Type -AssemblyName System.Speech
$text = {_powershell_quote(text)}
$outputPath = {_powershell_quote(output_path)}
$voiceCandidates = {ps_voice_array}
$synth = New-Object System.Speech.Synthesis.SpeechSynthesizer
try {{
    $selected = $null
    foreach ($voiceName in $voiceCandidates) {{
        try {{
            $synth.SelectVoice($voiceName)
            $selected = $voiceName
            break
        }} catch {{}}
    }}
    if (-not $selected) {{
        $fallbackVoice = $synth.GetInstalledVoices() |
            ForEach-Object {{ $_.VoiceInfo }} |
            Where-Object {{ $_.Culture.Name -like 'zh*' }} |
            Select-Object -First 1
        if ($fallbackVoice) {{
            $synth.SelectVoice($fallbackVoice.Name)
            $selected = $fallbackVoice.Name
        }}
    }}
    if (-not $selected) {{
        $fallbackVoice = $synth.GetInstalledVoices() |
            ForEach-Object {{ $_.VoiceInfo }} |
            Select-Object -First 1
        if ($fallbackVoice) {{
            $synth.SelectVoice($fallbackVoice.Name)
            $selected = $fallbackVoice.Name
        }}
    }}
    if (-not $selected) {{
        throw 'No Windows TTS voice is available.'
    }}

    $synth.Rate = 0
    $synth.Volume = 100
    $synth.SetOutputToWaveFile($outputPath)
    $synth.Speak($text)
    $synth.SetOutputToNull()
    Write-Output ("VOICE=" + $selected)
}} finally {{
    if ($synth) {{
        $synth.Dispose()
    }}
}}
"""

    try:
        result = subprocess.run(
            [
                "powershell.exe",
                "-NoProfile",
                "-NonInteractive",
                "-ExecutionPolicy",
                "Bypass",
                "-Command",
                ps_script,
            ],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="ignore",
            timeout=120,
        )
        if result.returncode != 0:
            print(f"Windows TTS failed: {result.stderr or result.stdout}")
            return False
        if not _audio_file_is_valid(output_path):
            print(f"Windows TTS created an invalid audio file: {output_path}")
            return False
        print(f"Windows TTS fallback generated audio: {output_path}")
        if result.stdout.strip():
            print(result.stdout.strip())
        return True
    except Exception as e:
        print(f"Windows TTS fallback error: {e}")
        return False

def call_tts(text, output_path, voice_id="longxia"):
    """Override legacy TTS logic to handle both JSON and direct audio responses."""
    try:
        tts_url = "https://dashscope.aliyuncs.com/api/v1/services/audio/tts/generation"
        headers = {
            "Authorization": f"Bearer {DASHSCOPE_API_KEY}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": "cosyvoice-tts-1.0",
            "input": {"text": text},
            "parameters": {
                "voice": voice_id,
                "format": "mp3",
                "sample_rate": 16000,
                "speed": 1.0,
                "volume": 1.0,
                "pitch": 1.0,
            },
        }

        response = requests.post(tts_url, headers=headers, json=payload, timeout=120)
        print(f"TTS request: {tts_url}")
        print(f"TTS status: {response.status_code}")
        content_type = response.headers.get("Content-Type", "").split(";")[0].strip().lower()
        print(f"TTS content type: {content_type}")

        if response.status_code == 200 and content_type.startswith("audio/"):
            if _save_audio_bytes(output_path, response.content):
                print(f"TTS audio saved: {output_path}")
                return True
            print(f"TTS direct audio was too small: {output_path}")
            return False

        response_text = response.text[:500]
        print(f"TTS response preview: {response_text}")

        if response.status_code == 200:
            result = response.json()
            audio_url = _extract_audio_url(result)
            if audio_url:
                print(f"TTS audio url: {audio_url}")
                audio_resp = requests.get(audio_url, timeout=60)
                if audio_resp.status_code == 200 and _save_audio_bytes(output_path, audio_resp.content):
                    print(f"TTS audio saved: {output_path}")
                    return True
                print(f"Download audio failed: status={audio_resp.status_code}, url={audio_url}")
        return False
    except Exception as e:
        print(f"TTS error: {e}")
        import traceback
        traceback.print_exc()
        return False

def generate_fallback_audio(text, output_path, voice_id="longxia"):
    if generate_windows_tts_audio(text, output_path, voice_id):
        return True
    return generate_dummy_wav(output_path)

CURRENT_AUDIO_FALLBACK_TEXT = ""
CURRENT_AUDIO_FALLBACK_VOICE_ID = "longxia"

def generate_dummy_audio(output_path):
    if _audio_file_is_valid(output_path):
        return True
    text = CURRENT_AUDIO_FALLBACK_TEXT or "这是当前页的讲解音频。"
    voice_id = CURRENT_AUDIO_FALLBACK_VOICE_ID or "longxia"
    return generate_fallback_audio(text, output_path, voice_id)

def render_slide_via_powershell(ppt_path, task_id):
    """使用PowerPoint COM导出幻灯片为图片"""
    import win32com.client
    import traceback
    import time

    image_urls = []
    ppt_app = None
    presentation = None
    
    try:
        print(f"开始PowerPoint渲染: {ppt_path}")
        print(f"文件存在: {os.path.exists(ppt_path)}")
        
        # 启动PowerPoint应用
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True  # 必须可见才能正确渲染
        print("PowerPoint应用已启动")
        
        # 等待应用启动
        time.sleep(1)

        # 打开PPT文件，使用更兼容的参数
        presentation = ppt_app.Presentations.Open(
            str(ppt_path), 
            ReadOnly=True,
            Untitled=False,
            WithWindow=True
        )
        print(f"PPT文件已打开，共有 {presentation.Slides.Count} 页")

        # 导出每一页幻灯片
        for slide_num in range(1, presentation.Slides.Count + 1):
            try:
                slide = presentation.Slides(slide_num)
                print(f"正在渲染第 {slide_num} 页")

                # 设置导出图片的尺寸（高分辨率）
                img_path = PPT_IMAGES_DIR / f"{task_id}_slide_{slide_num - 1}.png"
                
                # 使用更高分辨率导出
                slide.Export(
                    str(img_path), 
                    "PNG",
                    1280,  # 宽度
                    720   # 高度
                )
                
                # 检查文件是否成功创建
                if os.path.exists(img_path):
                    file_size = os.path.getsize(img_path)
                    print(f"已导出第 {slide_num} 页到 {img_path}, 大小: {file_size} 字节")
                    image_urls.append(f"/ppt_images/{task_id}_slide_{slide_num - 1}.png")
                else:
                    print(f"警告: 第 {slide_num} 页导出失败，文件不存在")
                    
            except Exception as slide_error:
                print(f"第 {slide_num} 页渲染错误: {slide_error}")
                traceback.print_exc()

        print(f"PowerPoint渲染完成，成功生成了 {len(image_urls)} 张图片")

    except Exception as e:
        print(f"PowerPoint渲染错误: {e}")
        traceback.print_exc()
        return []
    finally:
        # 确保正确关闭资源
        try:
            if presentation is not None:
                presentation.Close()
                print("演示文稿已关闭")
        except Exception as close_error:
            print(f"关闭演示文稿时出错: {close_error}")
            
        try:
            if ppt_app is not None:
                ppt_app.Quit()
                print("PowerPoint应用已关闭")
        except Exception as quit_error:
            print(f"关闭PowerPoint应用时出错: {quit_error}")
            
        # 等待一下确保进程完全退出
        time.sleep(0.5)

    return image_urls

def render_slide_to_image(slide, width=800, height=600):
    """将PPT幻灯片渲染为图片（备用方案）"""
    img = Image.new('RGB', (width, height), color='white')
    draw = ImageDraw.Draw(img)

    try:
        title_font = ImageFont.truetype("msyh.ttc", 28)
        content_font = ImageFont.truetype("msyh.ttc", 16)
    except:
        try:
            title_font = ImageFont.truetype("C:/Windows/Fonts/simhei.ttf", 28)
            content_font = ImageFont.truetype("C:/Windows/Fonts/simhei.ttf", 16)
        except:
            title_font = ImageFont.load_default()
            content_font = ImageFont.load_default()

    y_offset = 50
    title = slide.get('title', '')
    if title:
        draw.text((50, y_offset), title, fill='#333333', font=title_font)
        y_offset += 60

    content = slide.get('content', '')
    if content:
        lines = content.split('\n')
        for line in lines:
            if len(line) > 40:
                words = line.split()
                current_line = ""
                for word in words:
                    if len(current_line) + len(word) < 40:
                        current_line += word + " "
                    else:
                        draw.text((50, y_offset), current_line.strip(), fill='#666666', font=content_font)
                        y_offset += 25
                        current_line = word + " "
                if current_line:
                    draw.text((50, y_offset), current_line.strip(), fill='#666666', font=content_font)
                    y_offset += 25
            else:
                draw.text((50, y_offset), line, fill='#666666', font=content_font)
                y_offset += 25

    draw.text((width - 80, height - 40), f"第{slide['slide']}页", fill='#999999', font=content_font)

    return img

def parse_ppt_text(ppt_path):
    """解析PPT提取文字内容"""
    if not PPTX_AVAILABLE:
        return [{"slide": 1, "title": "PPT解析", "content": "python-pptx未安装，无法解析PPT"}]

    slides_content = []
    try:
        prs = Presentation(ppt_path)
        for idx, slide in enumerate(prs.slides):
            slide_text = {"slide": idx + 1, "title": "", "content": ""}

            # 提取标题和内容
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        if not slide_text["title"]:
                            slide_text["title"] = text[:100]
                        slide_text["content"] += text + "\n"

            if slide_text["content"].strip():
                slides_content.append(slide_text)

    except Exception as e:
        print(f"PPT解析错误: {e}")
        return [{"slide": 1, "title": "错误", "content": f"解析失败: {str(e)}"}]

    return slides_content

def generate_slide_images(ppt_path, slides, task_id):
    """为每一页幻灯片生成图片"""
    image_urls = []

    print("=" * 50)
    print("开始生成幻灯片图片")
    print(f"PPT路径: {ppt_path}")
    print(f"任务ID: {task_id}")
    print(f"幻灯片数量: {len(slides)}")
    print("=" * 50)

    # 首先尝试使用PowerPoint COM渲染 - 这是首选方案，必须优先执行
    powerpoint_success = False
    try:
        print("✓ 开始使用PowerPoint COM渲染完整PPT界面...")
        image_urls = render_slide_via_powershell(ppt_path, task_id)
        
        # 检查是否成功生成了图片
        if image_urls and len(image_urls) > 0:
            print(f"✓✓ PowerPoint COM渲染成功！生成了 {len(image_urls)} 张完整PPT图片")
            powerpoint_success = True
            
            # 确保返回的图片数量与幻灯片数量匹配或更多
            # 如果不匹配，为缺失的页面创建占位符
            if len(image_urls) < len(slides):
                print(f"⚠ 警告：生成的图片数量({len(image_urls)})少于幻灯片数量({len(slides)})")
                # 补充缺失的页面
                for idx in range(len(image_urls), len(slides)):
                    print(f"  为第 {idx+1} 页创建备用图片")
                    img = render_slide_to_image(slides[idx])
                    img_path = PPT_IMAGES_DIR / f"{task_id}_slide_{idx}.png"
                    img.save(str(img_path), "PNG")
                    image_urls.append(f"/ppt_images/{task_id}_slide_{idx}.png")
            
            return image_urls
        else:
            print("✗ PowerPoint COM返回空列表，继续尝试...")
    except ImportError as e:
        print(f"✗ win32com不可用: {e}")
        print("  请确保已安装pywin32: pip install pywin32")
    except Exception as e:
        print(f"✗ PowerPoint COM渲染异常: {e}")
        import traceback
        traceback.print_exc()

    # 如果PowerPoint COM完全失败，使用备用方案（仅文字）
    print("⚠ PowerPoint COM渲染失败，使用备用渲染方案（仅文字）...")
    print("  注意：备用方案只能显示文字，无法显示完整PPT界面")
    print("  建议：确保系统已安装Microsoft PowerPoint并尝试重新上传")
    
    image_urls = []
    for idx, slide in enumerate(slides):
        img = render_slide_to_image(slide)
        img_path = PPT_IMAGES_DIR / f"{task_id}_slide_{idx}.png"
        img.save(str(img_path), "PNG")
        image_urls.append(f"/ppt_images/{task_id}_slide_{idx}.png")
    
    print(f"备用方案生成了 {len(image_urls)} 张图片")
    return image_urls

@app.post("/api/upload_ppt")
async def upload_ppt(file: UploadFile = File(...)):
    """上传PPT文件并解析内容"""
    if not file.filename.endswith(('.ppt', '.pptx')):
        raise HTTPException(status_code=400, detail="只支持PPT/PPTX文件")

    task_id = str(uuid.uuid4())[:8]
    file_path = UPLOAD_DIR / f"{task_id}_{file.filename}"

    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # 解析PPT
    slides = parse_ppt_text(file_path)

    # 生成幻灯片图片
    slide_images = []
    if PPTX_AVAILABLE:
        slide_images = generate_slide_images(str(file_path), slides, task_id)

    return JSONResponse({
        "code": 0,
        "data": {
            "task_id": task_id,
            "filename": file.filename,
            "slides": slides,
            "slide_images": slide_images,
            "total_slides": len(slides)
        }
    })

@app.post("/api/generate_scripts")
async def generate_scripts(task_id: str):
    """为PPT每一页生成讲解脚本"""
    # 查找上传的PPT
    ppt_files = list(UPLOAD_DIR.glob(f"{task_id}_*"))
    if not ppt_files:
        raise HTTPException(status_code=404, detail="PPT文件不存在")

    ppt_path = ppt_files[0]
    slides = parse_ppt_text(ppt_path)

    scripts = []
    for slide in slides:
        script = call_llm_generate_script(slide.get('content', ''), slide['slide'])
        if script:
            scripts.append({
                "slide": slide['slide'],
                "title": slide.get('title', ''),
                "content": slide.get('content', ''),
                "script": script
            })
        else:
            scripts.append({
                "slide": slide['slide'],
                "title": slide.get('title', ''),
                "content": slide.get('content', ''),
                "script": slide.get('content', '')[:100]
            })

    # 保存脚本到文件
    scripts_path = TEMP_DIR / f"{task_id}_scripts.json"
    with open(scripts_path, 'w', encoding='utf-8') as f:
        json.dump(scripts, f, ensure_ascii=False, indent=2)

    return JSONResponse({
        "code": 0,
        "data": {
            "task_id": task_id,
            "scripts": scripts
        }
    })

@app.post("/api/generate_audio")
async def generate_audio(task_id: str, voice_id: str = "longxia"):
    """为讲解脚本生成语音"""
    scripts_path = TEMP_DIR / f"{task_id}_scripts.json"
    if not scripts_path.exists():
        raise HTTPException(status_code=404, detail="脚本文件不存在")

    with open(scripts_path, 'r', encoding='utf-8-sig') as f:
        scripts = json.load(f)

    audio_files = []
    for idx, script_data in enumerate(scripts):
        global CURRENT_AUDIO_FALLBACK_TEXT, CURRENT_AUDIO_FALLBACK_VOICE_ID
        CURRENT_AUDIO_FALLBACK_TEXT = str(script_data.get("script") or script_data.get("content") or "").strip()
        CURRENT_AUDIO_FALLBACK_VOICE_ID = voice_id
        audio_path = TEMP_DIR / f"{task_id}_audio_{idx}.mp3"
        audio_url = f"/temp/{task_id}_audio_{idx}.mp3"
        success = call_tts(
            str(script_data.get("script") or script_data.get("content") or "").strip(),
            str(audio_path),
            voice_id,
        )
        if not success:
            audio_path = TEMP_DIR / f"{task_id}_audio_{idx}.wav"
            audio_url = f"/temp/{task_id}_audio_{idx}.wav"
            success = generate_dummy_audio(str(audio_path))
            # TTS失败，使用备用音频方案
            success = generate_dummy_audio(str(audio_path))
        
        if success:
            audio_files.append({
                "slide": script_data['slide'],
                "audio_url": audio_url
            })
        else:
            audio_files.append({
                "slide": script_data['slide'],
                "audio_url": None
            })

    return JSONResponse({
        "code": 0,
        "data": {
            "task_id": task_id,
            "audios": audio_files
        }
    })

@app.get("/api/ppt_file/{task_id}")
async def get_ppt_file(task_id: str):
    """获取PPT文件用于预览"""
    ppt_files = list(UPLOAD_DIR.glob(f"{task_id}_*"))
    if not ppt_files:
        raise HTTPException(status_code=404, detail="PPT文件不存在")

    return FileResponse(str(ppt_files[0]), media_type="application/vnd.ms-powerpoint")

if __name__ == "__main__":
    import uvicorn
    print("=" * 50)
    print("PPT数字人视频制作服务启动中...")
    print(f"上传目录: {UPLOAD_DIR}")
    print(f"幻灯片图片目录: {PPT_IMAGES_DIR}")
    print("服务地址: http://localhost:8889")
    print("=" * 50)
    uvicorn.run(app, host="0.0.0.0", port=8889)
